"""
Пересобирает презентацию Directum RX:
- Слайд 1 — точная копия титульного оформления оригинала.
- Слайды 2–12 — единый шаблон в стиле слайда 3:
  фиолетовая полоска-акцент + заголовок Roboto Slab 32pt #3B1E5E + текст Roboto 18pt.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Палитра оригинала
PURPLE = RGBColor(0x3B, 0x1E, 0x5E)
MINT = RGBColor(0x4F, 0xD1, 0xB2)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xE8, 0xEE, 0xF4)
DARK = RGBColor(0x1F, 0x15, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Roboto Slab"
BODY_FONT = "Roboto"

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def inch(x):
    return Emu(int(x * 914400))


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    return tb, tf


def _set_para_bullet(p, char="•"):
    pPr = p._p.get_or_add_pPr()
    # remove existing bullets if any
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for child in pPr.findall(qn(tag)):
            pPr.remove(child)
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)
    pPr.set("marL", "285750")
    pPr.set("indent", "-285750")


def add_run(p, text, font=BODY_FONT, size=18, bold=False, color=DARK):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


# ---------- Слайд 1: титульник ----------
def build_title(slide):
    # Заголовок (центрированный)
    tb, tf = add_textbox(slide, inch(1.0), inch(1.23), inch(11.33), inch(2.6),
                         anchor=MSO_ANCHOR.MIDDLE)
    for i, line in enumerate(["Система электронного", "документооборота Directum RX"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        add_run(p, line, font=TITLE_FONT, size=44, bold=True, color=PURPLE)

    # TopLine — короткая фиолетовая черта сверху по центру
    add_rect(slide, inch(5.66), inch(1.5), inch(2.0), inch(0.042), PURPLE)

    # Метка ПРЕЗЕНТАЦИЯ
    tb, tf = add_textbox(slide, inch(4.33), inch(1.08), inch(4.66), inch(0.33))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "ПРЕЗЕНТАЦИЯ", font=BODY_FONT, size=14, bold=True, color=MINT)

    # Подзаголовок
    tb, tf = add_textbox(slide, inch(1.0), inch(3.94), inch(11.33), inch(1.8))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Современная СЭД для российского бизнеса",
            font=BODY_FONT, size=20, color=PURPLE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    add_run(p2, "Выполнил(а): ФИО студента, группа",
            font=BODY_FONT, size=16, color=SLATE)

    # BottomLine
    add_rect(slide, inch(4.33), inch(4.91), inch(4.66), inch(0.02), LIGHT)


# ---------- Контентный шаблон ----------
def build_content(slide, title, blocks, two_col=False):
    # TitleAccent (фиолетовая полоска слева)
    add_rect(slide, inch(0.5), inch(1.44), inch(0.75), inch(0.042), PURPLE)

    # Title
    tb, tf = add_textbox(slide, inch(0.92), inch(0.4), inch(11.5), inch(1.45))
    p = tf.paragraphs[0]
    add_run(p, title, font=TITLE_FONT, size=32, bold=True, color=PURPLE)

    if not two_col:
        # Body
        tb, tf = add_textbox(slide, inch(0.92), inch(2.0), inch(11.5), inch(4.76))
        first = True
        for block in blocks:
            kind = block["kind"]
            if kind == "subhead":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_before = Pt(0 if first else 10)
                p.space_after = Pt(4)
                add_run(p, block["text"], font=BODY_FONT, size=18, bold=True, color=PURPLE)
                first = False
            elif kind == "text":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_after = Pt(6)
                add_run(p, block["text"], font=BODY_FONT, size=18, color=DARK)
                first = False
            elif kind == "bullets":
                for item in block["items"]:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    p.space_after = Pt(4)
                    _set_para_bullet(p)
                    add_run(p, item, font=BODY_FONT, size=17, color=DARK)
                    first = False
    else:
        # Двухколоночный режим — blocks это [col_left, col_right]
        col_w = inch(5.55)
        positions = [inch(0.92), inch(6.86)]
        for col_blocks, x in zip(blocks, positions):
            tb, tf = add_textbox(slide, x, inch(2.0), col_w, inch(4.76))
            first = True
            for block in col_blocks:
                kind = block["kind"]
                if kind == "subhead":
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    p.space_before = Pt(0 if first else 10)
                    p.space_after = Pt(6)
                    add_run(p, block["text"], font=BODY_FONT, size=20, bold=True, color=PURPLE)
                    first = False
                elif kind == "bullets":
                    for item in block["items"]:
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        p.space_after = Pt(4)
                        _set_para_bullet(p)
                        add_run(p, item, font=BODY_FONT, size=17, color=DARK)
                        first = False


# ---------- Контент ----------
SLIDES = [
    None,  # титульник делается отдельно
    {
        "title": "Содержание",
        "blocks": [
            {"kind": "bullets", "items": [
                "01.  Понятие СЭД и её задачи",
                "02.  Преимущества внедрения",
                "03.  Общие сведения о Directum RX",
                "04.  Архитектура и платформы",
                "05.  Основные функции",
                "06.  Интеллектуальные сервисы Ario",
                "07.  Варианты поставки и тарифы",
                "08.  Безопасность",
                "09.  Примеры внедрения",
                "10.  Заключение",
            ]},
        ],
    },
    {
        "title": "Что такое СЭД?",
        "blocks": [
            {"kind": "subhead", "text": "Определение"},
            {"kind": "text", "text":
                "Система электронного документооборота (СЭД) — программный комплекс "
                "для автоматизации работы с электронными документами на всём жизненном "
                "цикле: от создания до архивного хранения."},
            {"kind": "subhead", "text": "Основные задачи СЭД:"},
            {"kind": "bullets", "items": [
                "создание, согласование и подписание документов",
                "регистрация, маршрутизация и контроль исполнения",
                "хранение и поиск, разграничение прав доступа",
                "обеспечение юридической значимости электронной подписью",
            ]},
        ],
    },
    {
        "title": "Преимущества внедрения СЭД",
        "blocks": [
            {"kind": "bullets", "items": [
                "Ускорение бизнес-процессов — сокращение времени согласования и подписания документов в 3–5 раз",
                "Снижение затрат — экономия на бумаге, печати, доставке и хранении документов",
                "Прозрачность и контроль — отслеживание статуса документов и истории всех действий",
                "Безопасность данных — разграничение прав доступа, защита от потери и несанкционированного использования",
                "Удалённая работа — доступ к документам из любой точки через веб и мобильные приложения",
            ]},
        ],
    },
    {
        "title": "Directum RX — общие сведения",
        "blocks": [
            {"kind": "bullets", "items": [
                "Разработчик: ООО «ДИРЕКТУМ» (Россия)",
                "Тип системы: ECM/СЭД с искусственным интеллектом",
                "Аудитория: средний и крупный бизнес, госсектор",
                "Реестр отечественного ПО: включена, подтверждённая импортонезависимость",
                "Распространение: более 3 500 компаний, до 50 000 пользователей",
            ]},
        ],
    },
    {
        "title": "Архитектура и платформы",
        "blocks": [
            {"kind": "bullets", "items": [
                "Микросервисная архитектура — гибкая, масштабируемая, легко настраиваемая экосистема",
                "Веб-клиент — полная работа с документами и задачами через браузер",
                "Мобильные приложения — Directum Solo и Jazz для iOS и Android",
                "Кроссплатформенность — поддержка Windows, Linux и macOS",
                "Интеграция — связь с 1С, ERP-системами, сервисами ЭДО (Synerdocs, Диадок)",
            ]},
        ],
    },
    {
        "title": "Основные функции Directum RX",
        "blocks": [
            {"kind": "bullets", "items": [
                "Управление документами — полный цикл, от создания и согласования до архивирования",
                "Управление процессами — маршруты согласования, поручения, контроль исполнения",
                "Делопроизводство — регистрация, рассмотрение корреспонденции, контроль исполнения",
                "Договоры — шаблоны, согласование, подписание ЭП, реестр обязательств",
                "Обмен с контрагентами — юридически значимый ЭДО через Synerdocs и Диадок",
                "Кадровые процессы (КЭДО) — приём, отпуска, командировки, удалённое подписание",
            ]},
        ],
    },
    {
        "title": "Интеллектуальные сервисы Directum Ario",
        "blocks": [
            {"kind": "bullets", "items": [
                "Компьютерное зрение — распознавание сканов и фотографий документов (OCR)",
                "Машинное обучение — автоматическая классификация документов по типам",
                "Извлечение реквизитов — автозаполнение карточек документов из текста",
                "Сравнение версий — поиск различий между редакциями договоров",
                "Поиск рисков — выявление проблемных условий в текстах договоров",
            ]},
        ],
    },
    {
        "title": "Варианты поставки и тарифы",
        "two_col": True,
        "blocks": [
            [  # левая колонка
                {"kind": "subhead", "text": "Локальная (On-Premise)"},
                {"kind": "bullets", "items": [
                    "Установка на серверах заказчика",
                    "Полный контроль данных",
                    "Тарифы: Basic, Enterprise, Intelligence",
                    "Подходит для крупных компаний",
                ]},
            ],
            [  # правая колонка
                {"kind": "subhead", "text": "Облачная (SaaS)"},
                {"kind": "bullets", "items": [
                    "Размещение в ЦОД на территории РФ",
                    "Быстрый старт, без капвложений",
                    "Тарифы: RX100 (до 200), DRX200",
                    "Подходит для среднего и малого бизнеса",
                ]},
            ],
        ],
    },
    {
        "title": "Безопасность и соответствие законодательству",
        "blocks": [
            {"kind": "bullets", "items": [
                "Электронная подпись — ПЭП и УКЭП, регламентируется 63-ФЗ «Об электронной подписи»",
                "Криптография — поддержка ГОСТ Р 34.10-2012; лицензия ФСБ России",
                "Персональные данные — соответствие 152-ФЗ «О персональных данных»",
                "Коммерческая тайна — соответствие 98-ФЗ «О коммерческой тайне»",
                "Контроль доступа — гибкая настройка прав, аудит действий, история работы с документами",
            ]},
        ],
    },
    {
        "title": "Примеры внедрения",
        "blocks": [
            {"kind": "bullets", "items": [
                "Росбанк — интеллектуальное распознавание документов от госорганов и клиентов",
                "Группа компаний «Эталон» — переход на единую корпоративную СЭД",
                "ESTEL — кадровый электронный документооборот (КЭДО)",
                "АО «Роскартография» — автоматизация делопроизводства и согласований",
                "ОДК — интеллектуальная обработка входящей корреспонденции",
                "Правительства регионов РФ — Тюменская, Омская, Курганская, Владимирская области",
            ]},
        ],
    },
    {
        "title": "Заключение",
        "blocks": [
            {"kind": "subhead", "text": "Универсальное решение"},
            {"kind": "text", "text": "Directum RX закрывает весь цикл документооборота — от создания документа до архива."},
            {"kind": "subhead", "text": "Российская разработка"},
            {"kind": "text", "text": "Полное соответствие законодательству РФ и независимость от зарубежного ПО."},
            {"kind": "subhead", "text": "Интеллектуальные сервисы"},
            {"kind": "text", "text": "ИИ Directum Ario автоматизирует рутинные операции и снижает нагрузку на сотрудников."},
            {"kind": "subhead", "text": "Гибкость поставки"},
            {"kind": "text", "text": "Облачные и локальные варианты подходят бизнесу любого масштаба — от 50 до 50 000 пользователей."},
            {"kind": "text", "text": "Источники: directum.ru, ru.wikipedia.org/wiki/Directum, tadviser.ru"},
        ],
    },
]


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Слайд 1 — титульник
    s1 = prs.slides.add_slide(blank)
    build_title(s1)

    # Слайды 2–12
    for spec in SLIDES[1:]:
        s = prs.slides.add_slide(blank)
        build_content(
            s,
            title=spec["title"],
            blocks=spec["blocks"],
            two_col=spec.get("two_col", False),
        )

    out = "Claude-for-Microsoft-365.pptx"
    prs.save(out)
    print(f"OK: {out}  slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
